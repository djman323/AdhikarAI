"""Add algorithms slides to existing PPTX presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Load existing presentation
prs = Presentation('Adhikar_AI_Stylish_Simple_Flow_HIL.pptx')

# Color scheme
BG_COLOR = RGBColor(7, 12, 22)  # #070C16
ACCENT_GOLD = RGBColor(231, 185, 107)  # #e7b96b
ACCENT_TEAL = RGBColor(0, 200, 200)  # #00c8c8
TEXT_COLOR = RGBColor(240, 242, 247)  # Light text
CARD_BG = RGBColor(20, 30, 50)  # Card background

def create_blank_slide_with_bg(prs):
    """Create blank slide with dark background."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    return slide

def add_title_to_slide(slide, title_text):
    """Add title to slide."""
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.LEFT

def add_content_box(slide, left, top, width, height, title, content_lines):
    """Add a content box with title and bullet points."""
    # Title
    title_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = False
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    
    # Content
    content_box = slide.shapes.add_textbox(left, top + Inches(0.55), width, height - Inches(0.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# ===== SLIDE 1: Algorithms Overview =====
slide1 = create_blank_slide_with_bg(prs)
add_title_to_slide(slide1, "Data Processing Algorithms")

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.paragraphs[0]
p.text = "12 complementary algorithms across 4 phases transform raw PDFs into reliable legal answers"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_COLOR
p.alignment = PP_ALIGN.LEFT

# Four phases - describe each
phases = [
    ("Phase 1: Document Loading", ["• PyPDFLoader: Extract text & metadata", "• RecursiveCharacterTextSplitter: Create parent/child chunks"]),
    ("Phase 2: Vectorization", ["• Embeddings: Convert chunks to 384-dim vectors", "• FAISS: Index for fast similarity search"]),
    ("Phase 3: Retrieval", ["• Dense Search: Semantic similarity (FAISS)", "• BM25: Keyword matching", "• RRF: Merge both signals", "• Cross-Encoder: Final ranking"]),
    ("Phase 4: Generation", ["• Prompt Assembly: Build final input", "• Qwen 2.5: Generate legal answer"])
]

col_width = Inches(4.2)
row_height = Inches(3.2)

positions = [
    (Inches(0.5), Inches(2.1)),
    (Inches(5.2), Inches(2.1)),
    (Inches(0.5), Inches(5.4)),
    (Inches(5.2), Inches(5.4))
]

for idx, (title, content) in enumerate(phases):
    left, top = positions[idx]
    add_content_box(slide1, left, top, col_width, row_height, title, content)


# ===== SLIDE 2: Hybrid Retrieval Pipeline =====
slide2 = create_blank_slide_with_bg(prs)
add_title_to_slide(slide2, "Hybrid Retrieval: Why It Works")

# Main explanation
main_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(1.2))
main_frame = main_box.text_frame
main_frame.word_wrap = True
p = main_frame.paragraphs[0]
p.text = "Constitutional law requires BOTH semantic understanding AND exact legal terminology. We combine 4 ranking signals:"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

# Four ranking approaches
ranking_data = [
    {
        "title": "1. Dense Search (FAISS)",
        "content": ["Semantic similarity", "Vector embeddings", "Catches meaning variations", "Example: 'right to life' → Article 21"]
    },
    {
        "title": "2. BM25 Keyword Search",
        "content": ["Exact term matching", "Word frequency scoring", "Preserves legal terms", "Example: 'habeas corpus' exact match"]
    },
    {
        "title": "3. RRF Fusion",
        "content": ["Merges both rankings", "Formula: 1/(60+rank)", "Avoids single-method bias", "Best of both worlds"]
    },
    {
        "title": "4. Cross-Encoder Reranking",
        "content": ["Fine-tuned legal ranker", "Legal relevance scoring", "Most accurate selection", "Final passage prioritization"]
    }
]

col_width = Inches(4.15)
row_height = Inches(2.8)

positions2 = [
    (Inches(0.5), Inches(2.8)),
    (Inches(5.15), Inches(2.8)),
    (Inches(0.5), Inches(5.7)),
    (Inches(5.15), Inches(5.7))
]

for idx, ranking in enumerate(ranking_data):
    left, top = positions2[idx]
    add_content_box(slide2, left, top, col_width, row_height, ranking['title'], ranking['content'])


# ===== SLIDE 3: Query Processing Flow =====
slide3 = create_blank_slide_with_bg(prs)
add_title_to_slide(slide3, "Query Processing: Step by Step")

# Example query
query_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.6))
query_frame = query_box.text_frame
p = query_frame.paragraphs[0]
p.text = "Example: 'What is my right if I am wrongfully detained?'"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_GOLD

# Steps
steps = [
    ("Step 1: Embed Question", "Convert 'wrongfully detained' to 384-dim vector using all-MiniLM-L6-v2"),
    ("Step 2: Dense Search", "FAISS finds 12 chunks near 'habeas corpus', 'Article 32', 'Personal Liberty'"),
    ("Step 3: BM25 Search", "Keyword search finds 12 chunks containing 'detention', 'writ', 'Article 21'"),
    ("Step 4: RRF Merge", "Combine both lists → ~15 unique chunks with fused scores"),
    ("Step 5: Cross-Encoder Rank", "Rank all 15 → Top 3 about habeas corpus and emergency protections"),
    ("Step 6: Add Context", "Expand top 3 child chunks with parent chunks for full Article context"),
    ("Step 7: Build Prompt", "Assemble: system + memory + retrieved context + question"),
    ("Step 8: Generate Answer", "Qwen 2.5 (temp=0.0) produces: 'Your right is protected under Articles 21-22...'"),
]

start_top = Inches(2.1)
line_height = Inches(0.42)

for idx, (step_title, step_desc) in enumerate(steps):
    top = start_top + (idx * line_height)
    
    # Step number circle
    step_num = idx + 1
    left = Inches(0.5)
    
    # Step text
    step_box = slide3.shapes.add_textbox(left + Inches(0.4), top, Inches(8.6), Inches(0.35))
    step_frame = step_box.text_frame
    step_frame.word_wrap = True
    
    p = step_frame.paragraphs[0]
    p.text = f"❱ {step_title}: {step_desc}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_COLOR
    p.space_before = Pt(1)
    p.space_after = Pt(1)


# ===== SLIDE 4: Key Formulas =====
slide4 = create_blank_slide_with_bg(prs)
add_title_to_slide(slide4, "Key Formulas & Parameters")

# RRF Formula
rrf_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
rrf_frame = rrf_box.text_frame
rrf_frame.word_wrap = True

p = rrf_frame.paragraphs[0]
p.text = "Reciprocal Rank Fusion (RRF)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

p = rrf_frame.add_paragraph()
p.text = "combined_score = 1/(60 + dense_rank) + 1/(60 + bm25_rank)"
p.font.size = Pt(12)
p.font.color.rgb = ACCENT_GOLD
p.level = 0

p = rrf_frame.add_paragraph()
p.text = "Why 60? Acts as smoothing constant. Prevents over-penalizing items ranked slightly lower."
p.font.size = Pt(10)
p.font.color.rgb = TEXT_COLOR
p.level = 0

# Chunking parameters
chunk_box = slide4.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.2), Inches(3))
chunk_frame = chunk_box.text_frame
chunk_frame.word_wrap = True

p = chunk_frame.paragraphs[0]
p.text = "Hierarchical Chunking"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

p = chunk_frame.add_paragraph()
p.text = "Parent Chunks"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = TEXT_COLOR
p.level = 0

p = chunk_frame.add_paragraph()
p.text = "Size: 2200 chars, Overlap: 220 chars\n→ Preserves broader context"
p.font.size = Pt(9)
p.font.color.rgb = TEXT_COLOR
p.level = 1

p = chunk_frame.add_paragraph()
p.text = "Child Chunks"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = TEXT_COLOR
p.level = 0

p = chunk_frame.add_paragraph()
p.text = "Size: 700 chars, Overlap: 100 chars\n→ Improves retrieval precision"
p.font.size = Pt(9)
p.font.color.rgb = TEXT_COLOR
p.level = 1

# Performance
perf_box = slide4.shapes.add_textbox(Inches(4.9), Inches(3.2), Inches(4.6), Inches(3))
perf_frame = perf_box.text_frame
perf_frame.word_wrap = True

p = perf_frame.paragraphs[0]
p.text = "Query Performance"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL

components = [
    ("Dense Search", "~10ms"),
    ("BM25 Search", "~5ms"),
    ("RRF Merge", "Negligible"),
    ("Cross-Encoder", "~50ms"),
    ("Qwen Generation", "~1-3s"),
    ("Total Latency", "~2-4s")
]

for comp, time in components:
    p = perf_frame.add_paragraph()
    p.text = f"{comp}: {time}"
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_COLOR
    p.level = 0

# Save presentation
prs.save('Adhikar_AI_Stylish_Simple_Flow_HIL.pptx')
print("✅ Added 4 algorithm slides to presentation!")
print("📊 Slides added:")
print("   - Slide 11: Data Processing Algorithms (4-phase overview)")
print("   - Slide 12: Hybrid Retrieval Pipeline (4 ranking signals)")
print("   - Slide 13: Query Processing Flow (8 steps with example)")
print("   - Slide 14: Key Formulas & Parameters (RRF, chunking, performance)")
